#!/usr/bin/env python3
"""PRD Markdown → PowerPoint 生成器

用法:
    python3 prd2pptx.py \
        --input <prd.md> \
        --output <output.pptx> \
        --outline <outline.json> \
        --color-scheme <scheme-name> \
        --style <style-name> \
        --font-zh <中文字体> \
        --font-en <英文字体>
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("错误: 需要安装 python-pptx。运行: pip3 install python-pptx", file=sys.stderr)
    sys.exit(1)

# ── 配色方案 ──────────────────────────────────────────────

COLOR_SCHEMES = {
    "business-authority": {
        "primary": "#1E3A5F", "secondary": "#2C5282", "accent": "#C9A84C",
        "background": "#FFFFFF", "text": "#1A202C",
    },
    "platinum-white-gold": {
        "primary": "#374151", "secondary": "#6B7280", "accent": "#D4AF37",
        "background": "#FAFAFA", "text": "#111827",
    },
    "tech-blue": {
        "primary": "#1D4ED8", "secondary": "#3B82F6", "accent": "#60A5FA",
        "background": "#F8FAFC", "text": "#1E293B",
    },
    "dreamy-creative": {
        "primary": "#7C3AED", "secondary": "#A78BFA", "accent": "#F472B6",
        "background": "#FAF5FF", "text": "#1F2937",
    },
    "bohemian": {
        "primary": "#DC2626", "secondary": "#F97316", "accent": "#FBBF24",
        "background": "#FFFBEB", "text": "#292524",
    },
    "art-food": {
        "primary": "#B45309", "secondary": "#D97706", "accent": "#92400E",
        "background": "#FEF3C7", "text": "#1C1917",
    },
    "coastal-coral": {
        "primary": "#0891B2", "secondary": "#06B6D4", "accent": "#FB923C",
        "background": "#F0FDFA", "text": "#164E63",
    },
    "nature-outdoor": {
        "primary": "#15803D", "secondary": "#22C55E", "accent": "#86EFAC",
        "background": "#F0FDF4", "text": "#14532D",
    },
    "forest-eco": {
        "primary": "#166534", "secondary": "#4ADE80", "accent": "#A3E635",
        "background": "#ECFDF5", "text": "#052E16",
    },
    "vintage-academic": {
        "primary": "#7C2D12", "secondary": "#9A3412", "accent": "#C2410C",
        "background": "#FFF7ED", "text": "#431407",
    },
    "education-chart": {
        "primary": "#1B4F72", "secondary": "#2980B9", "accent": "#F39C12",
        "background": "#FDF6EC", "text": "#1C2833",
    },
    "tech-vibrant": {
        "primary": "#4F46E5", "secondary": "#6366F1", "accent": "#EC4899",
        "background": "#EEF2FF", "text": "#1E1B4B",
    },
    "tech-nightscape": {
        "primary": "#0F172A", "secondary": "#1E293B", "accent": "#38BDF8",
        "background": "#020617", "text": "#E2E8F0",
    },
    "modern-health": {
        "primary": "#0D7377", "secondary": "#14B8A6", "accent": "#5EEAD4",
        "background": "#F0F9F4", "text": "#134E4A",
    },
    "artisan-handmade": {
        "primary": "#78350F", "secondary": "#A16207", "accent": "#CA8A04",
        "background": "#FEFCE8", "text": "#422006",
    },
    "elegant-fashion": {
        "primary": "#831843", "secondary": "#BE185D", "accent": "#F9A8D4",
        "background": "#FDF2F8", "text": "#500724",
    },
    "luxe-mystery": {
        "primary": "#1E1B4B", "secondary": "#312E81", "accent": "#C084FC",
        "background": "#0C0A1D", "text": "#DDD6FE",
    },
    "orange-mint": {
        "primary": "#EA580C", "secondary": "#F97316", "accent": "#34D399",
        "background": "#FFF7ED", "text": "#1C1917",
    },
}

# ── 风格定义 ──────────────────────────────────────────────

STYLES = {
    "corporate": {
        "title_size": 36, "body_size": 18, "line_spacing": 1.5,
        "title_align": PP_ALIGN.CENTER, "body_align": PP_ALIGN.LEFT,
    },
    "modern": {
        "title_size": 32, "body_size": 16, "line_spacing": 1.8,
        "title_align": PP_ALIGN.LEFT, "body_align": PP_ALIGN.LEFT,
    },
    "classic": {
        "title_size": 28, "body_size": 14, "line_spacing": 1.4,
        "title_align": PP_ALIGN.CENTER, "body_align": PP_ALIGN.LEFT,
    },
    "creative": {
        "title_size": 40, "body_size": 18, "line_spacing": 1.6,
        "title_align": PP_ALIGN.LEFT, "body_align": PP_ALIGN.LEFT,
    },
}

# ── 幻灯片尺寸常量 ───────────────────────────────────────

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.8)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN


# ── 工具函数 ──────────────────────────────────────────────

def hex_to_rgb(hex_str: str) -> RGBColor:
    """将 #RRGGBB 转换为 RGBColor"""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_font(run, size_pt: int, color_hex: str, font_zh: str, font_en: str, bold: bool = False):
    """设置 run 的字体属性"""
    run.font.size = Pt(size_pt)
    run.font.color.rgb = hex_to_rgb(color_hex)
    run.font.bold = bold
    run.font.name = font_en
    # 设置东亚字体
    from pptx.oxml.ns import qn
    rPr = run._r.get_or_add_rPr()
    ea = rPr.makeelement(qn("a:ea"), {"typeface": font_zh})
    # 移除已有的 ea 元素（如果有）
    for old_ea in rPr.findall(qn("a:ea")):
        rPr.remove(old_ea)
    rPr.append(ea)


def add_background(slide, color_hex: str):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color_hex)


def add_rect_shape(slide, left, top, width, height, fill_hex: str):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_hex)
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text: str, size_pt: int,
                 color_hex: str, font_zh: str, font_en: str,
                 bold: bool = False, alignment=PP_ALIGN.LEFT):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    set_font(run, size_pt, color_hex, font_zh, font_en, bold)
    return txBox


# ── 页面生成器 ────────────────────────────────────────────

def build_cover(prs, slide_data: dict, colors: dict, style: dict,
                font_zh: str, font_en: str):
    """生成封面页"""
    layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(layout)

    # 背景
    add_background(slide, colors["primary"])

    # 标题
    title_top = Inches(2.2)
    add_text_box(
        slide, MARGIN, title_top, CONTENT_WIDTH, Inches(1.5),
        slide_data["title"], style["title_size"] + 8,
        "#FFFFFF", font_zh, font_en, bold=True, alignment=PP_ALIGN.CENTER,
    )

    # 副标题
    if slide_data.get("bullets"):
        subtitle = slide_data["bullets"][0] if slide_data["bullets"] else ""
        add_text_box(
            slide, MARGIN, Inches(4.0), CONTENT_WIDTH, Inches(1.0),
            subtitle, style["body_size"] + 2,
            colors.get("accent", "#FFFFFF"), font_zh, font_en,
            alignment=PP_ALIGN.CENTER,
        )

    return slide


def build_section(prs, slide_data: dict, colors: dict, style: dict,
                  font_zh: str, font_en: str):
    """生成章节分隔页"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    add_background(slide, colors["background"])

    # 左侧色块装饰
    add_rect_shape(slide, Inches(0), Inches(0), Inches(0.4), SLIDE_HEIGHT, colors["primary"])

    # 章节标题
    add_text_box(
        slide, Inches(1.5), Inches(2.5), Inches(10), Inches(2.0),
        slide_data["title"], style["title_size"] + 4,
        colors["primary"], font_zh, font_en, bold=True,
        alignment=style["title_align"],
    )

    return slide


def build_content_bullet_list(prs, slide_data: dict, colors: dict, style: dict,
                              font_zh: str, font_en: str):
    """生成内容页 — 要点列表"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_background(slide, colors["background"])

    # 标题
    add_text_box(
        slide, MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(1.0),
        slide_data["title"], style["title_size"],
        colors["primary"], font_zh, font_en, bold=True,
        alignment=style["title_align"],
    )

    # 底部色条
    add_rect_shape(
        slide, MARGIN, Inches(1.3), Inches(2.0), Inches(0.05), colors["accent"]
    )

    # 要点
    bullet_top = Inches(1.8)
    for i, bullet in enumerate(slide_data.get("bullets", [])):
        y = bullet_top + Inches(i * 0.8)
        # 圆点
        add_text_box(
            slide, Inches(1.0), y, Inches(0.3), Inches(0.5),
            "\u2022", style["body_size"],
            colors["accent"], font_zh, font_en,
        )
        # 文字
        add_text_box(
            slide, Inches(1.4), y, Inches(10.5), Inches(0.6),
            bullet, style["body_size"],
            colors["text"], font_zh, font_en,
            alignment=style["body_align"],
        )

    return slide


def build_content_two_column(prs, slide_data: dict, colors: dict, style: dict,
                             font_zh: str, font_en: str):
    """生成内容页 — 双栏对比"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_background(slide, colors["background"])

    # 标题
    add_text_box(
        slide, MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(1.0),
        slide_data["title"], style["title_size"],
        colors["primary"], font_zh, font_en, bold=True,
        alignment=style["title_align"],
    )

    add_rect_shape(
        slide, MARGIN, Inches(1.3), Inches(2.0), Inches(0.05), colors["accent"]
    )

    bullets = slide_data.get("bullets", [])
    mid = len(bullets) // 2 if len(bullets) > 1 else len(bullets)
    left_bullets = bullets[:mid]
    right_bullets = bullets[mid:]

    col_width = Inches(5.5)
    for col_idx, col_bullets in enumerate([left_bullets, right_bullets]):
        x_offset = MARGIN + Inches(col_idx * 6.0)
        for i, bullet in enumerate(col_bullets):
            y = Inches(1.8) + Inches(i * 0.8)
            add_text_box(
                slide, x_offset, y, Inches(0.3), Inches(0.5),
                "\u2022", style["body_size"],
                colors["accent"], font_zh, font_en,
            )
            add_text_box(
                slide, x_offset + Inches(0.4), y, col_width, Inches(0.6),
                bullet, style["body_size"],
                colors["text"], font_zh, font_en,
            )

    return slide


def build_content_image_text(prs, slide_data: dict, colors: dict, style: dict,
                             font_zh: str, font_en: str):
    """生成内容页 — 图文混排（占位图区域）"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_background(slide, colors["background"])

    # 标题
    add_text_box(
        slide, MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(1.0),
        slide_data["title"], style["title_size"],
        colors["primary"], font_zh, font_en, bold=True,
        alignment=style["title_align"],
    )

    # 左侧占位图区域
    img_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        MARGIN, Inches(1.8), Inches(5.0), Inches(4.5),
    )
    img_shape.fill.solid()
    img_shape.fill.fore_color.rgb = hex_to_rgb(colors["secondary"])
    img_shape.fill.fore_color.brightness = 0.8
    img_shape.line.color.rgb = hex_to_rgb(colors["accent"])
    img_shape.line.width = Pt(2)

    # 占位提示
    add_text_box(
        slide, Inches(1.5), Inches(3.5), Inches(3.5), Inches(1.0),
        "[ 在此插入图片 ]", style["body_size"] - 2,
        colors["secondary"], font_zh, font_en, alignment=PP_ALIGN.CENTER,
    )

    # 右侧文字
    text_x = Inches(6.5)
    for i, bullet in enumerate(slide_data.get("bullets", [])):
        y = Inches(1.8) + Inches(i * 0.8)
        add_text_box(
            slide, text_x, y, Inches(5.5), Inches(0.6),
            f"\u2022  {bullet}", style["body_size"],
            colors["text"], font_zh, font_en,
        )

    return slide


def build_chart(prs, slide_data: dict, colors: dict, style: dict,
                font_zh: str, font_en: str):
    """生成数据图表页（占位区域）"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_background(slide, colors["background"])

    # 标题
    add_text_box(
        slide, MARGIN, Inches(0.5), CONTENT_WIDTH, Inches(1.0),
        slide_data["title"], style["title_size"],
        colors["primary"], font_zh, font_en, bold=True,
        alignment=style["title_align"],
    )

    # 图表占位区域
    chart_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(1.8), Inches(10.0), Inches(4.5),
    )
    chart_shape.fill.background()
    chart_shape.line.color.rgb = hex_to_rgb(colors["secondary"])
    chart_shape.line.width = Pt(1.5)
    chart_shape.line.dash_style = 2  # dash

    add_text_box(
        slide, Inches(4.5), Inches(3.5), Inches(4.0), Inches(1.0),
        "[ 在此插入图表 ]", style["body_size"],
        colors["secondary"], font_zh, font_en, alignment=PP_ALIGN.CENTER,
    )

    # 数据说明
    if slide_data.get("bullets"):
        for i, bullet in enumerate(slide_data["bullets"][:3]):
            y = Inches(6.5) + Inches(i * 0.3)
            add_text_box(
                slide, MARGIN, y, CONTENT_WIDTH, Inches(0.3),
                f"  {bullet}", style["body_size"] - 4,
                colors["text"], font_zh, font_en,
            )

    return slide


def build_end(prs, slide_data: dict, colors: dict, style: dict,
              font_zh: str, font_en: str):
    """生成结束页"""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_background(slide, colors["primary"])

    # 感谢语
    add_text_box(
        slide, MARGIN, Inches(2.5), CONTENT_WIDTH, Inches(1.5),
        slide_data.get("title", "谢谢"), style["title_size"] + 8,
        "#FFFFFF", font_zh, font_en, bold=True, alignment=PP_ALIGN.CENTER,
    )

    # 联系信息
    if slide_data.get("bullets"):
        contact = " | ".join(slide_data["bullets"])
        add_text_box(
            slide, MARGIN, Inches(4.5), CONTENT_WIDTH, Inches(1.0),
            contact, style["body_size"],
            colors.get("accent", "#FFFFFF"), font_zh, font_en,
            alignment=PP_ALIGN.CENTER,
        )

    return slide


# ── Content 页子布局循环分配 ──────────────────────────────

CONTENT_BUILDERS = [
    ("bullet-list", build_content_bullet_list),
    ("two-column", build_content_two_column),
    ("image-text", build_content_image_text),
]


def get_content_builder(sub_layout: str | None, auto_index: int):
    """根据 sub_layout 或自动循环选择 Content 子布局"""
    if sub_layout:
        for name, builder in CONTENT_BUILDERS:
            if name == sub_layout:
                return builder
    # 自动循环
    return CONTENT_BUILDERS[auto_index % len(CONTENT_BUILDERS)][1]


# ── 页面类型路由 ──────────────────────────────────────────

PAGE_BUILDERS = {
    "cover": build_cover,
    "section": build_section,
    "chart": build_chart,
    "end": build_end,
}


# ── 主函数 ────────────────────────────────────────────────

def generate(outline_path: str, output_path: str, color_scheme: str,
             style_name: str, font_zh: str, font_en: str):
    """生成 PPTX"""
    # 加载大纲
    with open(outline_path, "r", encoding="utf-8") as f:
        outline = json.load(f)

    # 获取配色和风格
    colors = COLOR_SCHEMES.get(color_scheme)
    if not colors:
        print(f"错误: 未知配色方案 '{color_scheme}'。可用: {', '.join(COLOR_SCHEMES.keys())}", file=sys.stderr)
        sys.exit(1)

    style = STYLES.get(style_name)
    if not style:
        print(f"错误: 未知风格 '{style_name}'。可用: {', '.join(STYLES.keys())}", file=sys.stderr)
        sys.exit(1)

    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    content_auto_index = 0

    for slide_data in outline:
        page_type = slide_data.get("page_type", "content")

        if page_type == "content":
            builder = get_content_builder(
                slide_data.get("sub_layout"), content_auto_index
            )
            builder(prs, slide_data, colors, style, font_zh, font_en)
            content_auto_index += 1
        elif page_type in PAGE_BUILDERS:
            PAGE_BUILDERS[page_type](prs, slide_data, colors, style, font_zh, font_en)
        else:
            print(f"警告: 未知页面类型 '{page_type}'，降级为 bullet-list", file=sys.stderr)
            build_content_bullet_list(prs, slide_data, colors, style, font_zh, font_en)

    # 保存
    Path(output_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)

    result = {
        "path": str(Path(output_path).resolve()),
        "slide_count": len(prs.slides),
        "color_scheme": color_scheme,
        "style": style_name,
    }
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return result


def main():
    parser = argparse.ArgumentParser(description="PRD → PPTX 生成器")
    parser.add_argument("--input", required=False, help="PRD Markdown 文件路径（当前未使用，保留）")
    parser.add_argument("--output", required=True, help="输出 .pptx 文件路径")
    parser.add_argument("--outline", required=True, help="大纲 JSON 文件路径")
    parser.add_argument("--color-scheme", default="tech-blue", help="配色方案名称")
    parser.add_argument("--style", default="modern", help="风格名称")
    parser.add_argument("--font-zh", default="PingFang SC", help="中文字体")
    parser.add_argument("--font-en", default="Helvetica Neue", help="英文字体")
    args = parser.parse_args()

    generate(
        outline_path=args.outline,
        output_path=args.output,
        color_scheme=args.color_scheme,
        style_name=args.style,
        font_zh=args.font_zh,
        font_en=args.font_en,
    )


if __name__ == "__main__":
    main()
