#!/usr/bin/env python3
"""PPT QA 验证脚本

用法:
    python3 verify.py <output.pptx> [--outline <outline.json>]

输出 JSON 格式的验证报告到 stdout。
"""

import argparse
import json
import re
import sys
from pathlib import Path

# ── 占位符检测模式 ────────────────────────────────────────

PLACEHOLDER_PATTERNS = [
    re.compile(r"\{[^}]{2,}\}"),         # {变量名}
    re.compile(r"\[.*?占位.*?\]"),        # [xxx占位]
    re.compile(r"\bTODO\b"),             # TODO
    re.compile(r"\bTBD\b"),              # TBD
    re.compile(r"Lorem ipsum", re.I),    # Lorem ipsum
]

# 允许的占位文本（图表/图片区域的提示不算错误）
ALLOWED_PLACEHOLDERS = [
    "在此插入图表",
    "在此插入图片",
]


def extract_text_markitdown(pptx_path: str) -> str | None:
    """尝试用 markitdown 提取文本"""
    try:
        from markitdown import MarkItDown
        m = MarkItDown()
        result = m.convert(pptx_path)
        return result.text_content
    except ImportError:
        return None
    except Exception as e:
        print(f"markitdown 提取失败: {e}", file=sys.stderr)
        return None


def extract_text_pptx(pptx_path: str) -> tuple[str, int]:
    """用 python-pptx 直接提取文本，返回 (全文, 幻灯片数)"""
    try:
        from pptx import Presentation
    except ImportError:
        print("错误: 需要安装 python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(pptx_path)
    slide_count = len(prs.slides)
    texts = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        texts.append("\n".join(slide_texts))

    return "\n---\n".join(texts), slide_count


def check_placeholders(text: str) -> list[str]:
    """检测占位符残留"""
    found = []
    for pattern in PLACEHOLDER_PATTERNS:
        for match in pattern.finditer(text):
            matched_text = match.group()
            # 跳过允许的占位文本
            if any(allowed in matched_text for allowed in ALLOWED_PLACEHOLDERS):
                continue
            found.append(matched_text)
    return found


def check_outline_coverage(text: str, outline: list[dict]) -> dict:
    """检查大纲覆盖率"""
    total_titles = 0
    matched_titles = 0
    total_bullets = 0
    matched_bullets = 0

    for slide_data in outline:
        title = slide_data.get("title", "")
        if title:
            total_titles += 1
            if title in text:
                matched_titles += 1

        for bullet in slide_data.get("bullets", []):
            if bullet:
                total_bullets += 1
                if bullet in text:
                    matched_bullets += 1

    title_coverage = matched_titles / total_titles if total_titles > 0 else 1.0
    bullet_coverage = matched_bullets / total_bullets if total_bullets > 0 else 1.0

    return {
        "title_coverage": round(title_coverage, 2),
        "bullet_coverage": round(bullet_coverage, 2),
        "title_detail": f"{matched_titles}/{total_titles}",
        "bullet_detail": f"{matched_bullets}/{total_bullets}",
    }


def verify(pptx_path: str, outline_path: str | None = None) -> dict:
    """执行验证流程"""
    pptx_file = Path(pptx_path)
    if not pptx_file.exists():
        return {
            "status": "fail",
            "errors": [f"文件不存在: {pptx_path}"],
            "warnings": [],
        }

    # 文件大小检查
    file_size_kb = pptx_file.stat().st_size / 1024
    warnings = []
    errors = []

    if file_size_kb > 1024:
        warnings.append(f"文件偏大: {file_size_kb:.0f}KB")

    # 提取文本
    text_md = extract_text_markitdown(pptx_path)
    text_pptx, slide_count = extract_text_pptx(pptx_path)
    text = text_md if text_md else text_pptx

    extraction_method = "markitdown" if text_md else "python-pptx"

    # 占位符检测
    placeholders = check_placeholders(text)
    placeholder_free = len(placeholders) == 0
    if not placeholder_free:
        errors.append(f"发现占位符残留: {placeholders}")

    # 大纲覆盖率（如果提供了 outline）
    coverage = {}
    if outline_path:
        outline_file = Path(outline_path)
        if outline_file.exists():
            with open(outline_file, "r", encoding="utf-8") as f:
                outline = json.load(f)
            coverage = check_outline_coverage(text, outline)

            # 页数匹配
            outline_count = len(outline)
            slide_count_match = slide_count == outline_count
            if not slide_count_match:
                warnings.append(f"页数不匹配: 大纲 {outline_count} 页 vs 实际 {slide_count} 页")

            # 覆盖率低警告
            if coverage.get("bullet_coverage", 1.0) < 0.9:
                warnings.append(f"要点覆盖率偏低: {coverage['bullet_coverage']}")
        else:
            warnings.append(f"大纲文件不存在: {outline_path}")

    # 判断状态
    if errors:
        status = "fail"
    elif warnings:
        status = "warn"
    else:
        status = "pass"

    report = {
        "status": status,
        "slide_count": slide_count,
        "file_size_kb": round(file_size_kb, 1),
        "extraction_method": extraction_method,
        "checks": {
            "placeholder_free": placeholder_free,
            **coverage,
        },
        "warnings": warnings,
        "errors": errors,
    }

    return report


def main():
    parser = argparse.ArgumentParser(description="PPT QA 验证")
    parser.add_argument("pptx", help=".pptx 文件路径")
    parser.add_argument("--outline", help="大纲 JSON 文件路径（可选）")
    args = parser.parse_args()

    report = verify(args.pptx, args.outline)
    print(json.dumps(report, ensure_ascii=False, indent=2))

    # 返回码
    sys.exit(0 if report["status"] in ("pass", "warn") else 1)


if __name__ == "__main__":
    main()
